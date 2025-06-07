# extrator_generico.py
import sys
from pptx import Presentation

def extrair_texto_de_pptx(caminho_do_arquivo):
    """
    Abre um arquivo .pptx e extrai todo o texto de cada slide.

    Args:
        caminho_do_arquivo (str): O caminho para o arquivo .pptx.

    Returns:
        dict: Um dicionário onde a chave é o número do slide e o valor é
              uma lista de textos encontrados naquele slide.
              Retorna None se o arquivo não for encontrado.
    """
    try:
        # Abre o arquivo de apresentação
        apresentacao = Presentation(caminho_do_arquivo)
        
        # Dicionário para armazenar o texto extraído
        dados_dos_slides = {}

        # Itera sobre cada slide na apresentação
        for i, slide in enumerate(apresentacao.slides):
            numero_slide = i + 1
            dados_dos_slides[numero_slide] = []
            
            # Itera sobre cada forma (caixa de texto, título, etc.) no slide
            for shape in slide.shapes:
                # Verifica se a forma contém texto
                if not shape.has_text_frame:
                    continue
                
                # Extrai o texto da forma
                texto_da_forma = shape.text.strip()
                if texto_da_forma:
                    dados_dos_slides[numero_slide].append(texto_da_forma)
        
        return dados_dos_slides

    except FileNotFoundError:
        print(f"🚨 ERRO: O arquivo '{caminho_do_arquivo}' não foi encontrado.")
        return None
    except Exception as e:
        print(f"Ocorreu um erro ao processar o arquivo: {e}")
        return None

def imprimir_dados_extraidos(dados):
    """Imprime os dados extraídos de forma legível."""
    if not dados:
        print("Nenhum texto foi extraído.")
        return
        
    print("--- INÍCIO DA EXTRAÇÃO DO PPTX ---")
    for numero_slide, textos in dados.items():
        print(f"\n--- SLIDE {numero_slide} ---")
        if textos:
            for texto in textos:
                # Imprime o texto, substituindo quebras de linha internas por um espaço
                print(texto.replace('\n', ' '))
        else:
            print("(Nenhum texto encontrado neste slide)")
    print("\n--- FIM DA EXTRAÇÃO ---")


if __name__ == '__main__':
    # Verifica se o usuário forneceu o caminho do arquivo como argumento
    if len(sys.argv) < 2:
        print("Uso: python extrator_generico.py <caminho_para_seu_arquivo.pptx>")
        sys.exit(1) # Sai do script se o argumento estiver faltando

    caminho_do_arquivo_pptx = sys.argv[1]
    
    dados_extraidos = extrair_texto_de_pptx(caminho_do_arquivo_pptx)
    
    if dados_extraidos:
        imprimir_dados_extraidos(dados_extraidos)