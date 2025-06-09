# scripts/extrair_info_ppxt.py (Refatorado para ser importável)
import sys
from pptx import Presentation

def extrair_texto_de_pptx(caminho_do_arquivo_pptx, caminho_do_arquivo_saida):
    """
    Abre um arquivo .pptx, extrai o texto de cada slide e salva em um arquivo de texto.
    """
    try:
        apresentacao = Presentation(caminho_do_arquivo_pptx)
        
        with open(caminho_do_arquivo_saida, 'w', encoding='utf-8') as f_out:
            f_out.write("--- INÍCIO DA EXTRAÇÃO DO PPTX ---\n")
            for i, slide in enumerate(apresentacao.slides):
                f_out.write(f"\n--- SLIDE {i + 1} ---\n")
                textos_do_slide = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        texto_da_forma = shape.text.strip()
                        if texto_da_forma:
                            textos_do_slide.append(texto_da_forma.replace('\n', ' '))
                
                if textos_do_slide:
                    f_out.write("\n".join(textos_do_slide) + "\n")
                else:
                    f_out.write("(Nenhum texto encontrado neste slide)\n")
            f_out.write("\n--- FIM DA EXTRAÇÃO ---")
        return True

    except Exception as e:
        print(f"Ocorreu um erro ao processar o arquivo PPTX: {e}")
        return False

if __name__ == '__main__':
    if len(sys.argv) < 3:
        print("Uso: python extrair_info_ppxt.py <caminho_para_pptx> <caminho_para_txt_saida>")
        sys.exit(1)
    
    extrair_texto_de_pptx(sys.argv[1], sys.argv[2])